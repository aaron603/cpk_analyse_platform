"""
gen_ppt.py — 产线数据分析AI平台 架构介绍 PPT（基于公司模板）
运行：
  "C:\\Program Files\\Python314\\python.exe" gen_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 公司模板路径 ──────────────────────────────────────────────────────────────
TEMPLATE = (r'D:\恒湾科技\项目\Sakura\Apple ORU6228 N25N66'
            r'\N25N66 Production solution PA1.pptx')

# ── 公司配色（取自模板主题）──────────────────────────────────────────────────
C_RED    = RGBColor(0xC0, 0x00, 0x00)   # 深红 — 标题栏/卡片头
C_RED2   = RGBColor(0xFF, 0x00, 0x00)   # 亮红 — 强调/流程步骤
C_DARK   = RGBColor(0x38, 0x08, 0x00)   # 深红褐 — 次要标题
C_PINK   = RGBColor(0xFD, 0xE2, 0xE1)   # 浅粉 — 内容底色
C_CREAM  = RGBColor(0xF5, 0xEC, 0xE9)   # 极浅米 — 备用底色
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_GRAY   = RGBColor(0x33, 0x33, 0x33)
# 保留少量绿/橙 区分"已实现"/"待实现"状态
C_GREEN  = RGBColor(0x37, 0x60, 0x34)
C_ORANGE = RGBColor(0xBF, 0x6A, 0x02)

W = Inches(13.33)
H = Inches(7.5)


# ── 基础 ──────────────────────────────────────────────────────────────────────

def init_prs() -> Presentation:
    """打开模板，清空现有幻灯片，返回干净的 Presentation（保留 master/logo）。"""
    from pptx.oxml.ns import qn as _qn
    prs = Presentation(TEMPLATE)
    part = prs.part
    sldIdLst = part._element.find(_qn('p:sldIdLst'))

    # Collect rIds before removing
    rIds = [el.get(_qn('r:id')) for el in list(sldIdLst)]

    # Remove all entries from sldIdLst
    for el in list(sldIdLst):
        sldIdLst.remove(el)

    # Drop relationships so orphaned slide parts are not written on save
    for rId in rIds:
        part.rels._rels.pop(rId, None)

    return prs


def blank_slide(prs: Presentation):
    """添加空白幻灯片（layout 0 = 空白，自动继承 master logo）。"""
    return prs.slides.add_slide(prs.slide_layouts[0])


# ── 绘图工具 ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(13), bold=False, color=C_BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def title_bar(slide, title, subtitle=''):
    """红色顶部标题栏（留出右上角 Logo 区域）。"""
    add_rect(slide, Inches(0), Inches(0), Inches(11.3), Inches(1.3), fill=C_RED)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(10.8), Inches(0.75),
             size=Pt(26), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(10.8), Inches(0.38),
                 size=Pt(12), color=RGBColor(0xFF, 0xCC, 0xCC))


def footer(slide):
    add_rect(slide, Inches(0), Inches(7.22), W, Inches(0.28), fill=C_RED)
    add_text(slide, 'Zillnk Efficiency Improvement Group  |  产线数据分析AI平台 v2.0',
             Inches(0.3), Inches(7.22), Inches(12.7), Inches(0.28),
             size=Pt(9), color=C_WHITE, align=PP_ALIGN.CENTER)


def bullets(slide, items, l, t, w, h,
            size=Pt(12), color=C_BLACK, prefix='• '):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f'{prefix}{item}'
        run.font.size = size
        run.font.color.rgb = color


def card(slide, title, body_lines, l, t, w, h,
         title_bg=C_RED, body_bg=C_PINK,
         title_sz=Pt(12), body_sz=Pt(10.5)):
    th = Inches(0.36)
    add_rect(slide, l, t, w, th, fill=title_bg)
    add_text(slide, title,
             l + Inches(0.08), t + Pt(3), w - Inches(0.1), th,
             size=title_sz, bold=True, color=C_WHITE)
    add_rect(slide, l, t + th, w, h - th,
             fill=body_bg, line=title_bg, line_w=Pt(0.75))
    bullets(slide, body_lines,
            l + Inches(0.1), t + th + Pt(3),
            w - Inches(0.15), h - th - Pt(5),
            size=body_sz)


# ═══════════════════════════════════════════════════════════════════════════════
# Slides
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = blank_slide(prs)

    # Full red background (title page)
    add_rect(slide, Inches(0), Inches(0), W, H, fill=C_RED)
    # White band
    add_rect(slide, Inches(0), Inches(4.2), W, Inches(3.3), fill=C_WHITE)
    # Accent stripe
    add_rect(slide, Inches(0), Inches(4.1), W, Inches(0.12), fill=C_DARK)

    add_text(slide, '产线数据分析AI平台',
             Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.2),
             size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 'Zillnk Efficiency Improvement Group',
             Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.6),
             size=Pt(20), color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    add_text(slide, 'v2.0  |  功能架构说明',
             Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.5),
             size=Pt(15), color=RGBColor(0xFF, 0xE5, 0xE5), align=PP_ALIGN.CENTER)

    add_text(slide, '本地测试站数据分析  ·  CPK 计算  ·  故障定位关系库',
             Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.5),
             size=Pt(16), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    add_text(slide, 'Python 3  ·  tkinter GUI  ·  openpyxl  ·  SQLite  ·  Ollama LLM',
             Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.45),
             size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = blank_slide(prs)
    title_bar(slide, '平台概览', '设计目标 · 用户群体 · 三大功能模块')
    footer(slide)

    card(slide, '平台定位',
         ['供 Zillnk 质量工程师使用的本地数据分析工具',
          '无需联网，数据不出厂，保障数据安全',
          '支持 PA 板（B3/B40）、RRU（Apricot）等多产品',
          '配置持久化，重启自动恢复上次设置',
          '分析结果自动打开浏览器查看 HTML 报告'],
         Inches(0.3), Inches(1.45), Inches(3.9), Inches(2.5))

    card(slide, '功能一：本地测试站数据分析（已实现）',
         ['5种分析模式：最后pass / 所选文件夹 / 全部成功 / 含失败 / 仅失败',
          '自动遍历工站目录，提取 xlsx / json 测试文件',
          'CPK 计算 + 综合 HTML 报告（6 Tab）',
          '失败分析报告 + 故障关系库（SQLite）',
          '工站合并配置，多机台数据合并分析'],
         Inches(4.4), Inches(1.45), Inches(4.4), Inches(2.5),
         title_bg=C_GREEN, body_bg=RGBColor(0xE8,0xF5,0xE9))

    card(slide, '功能二 / 三（待实现）',
         ['功能二：深科技 MES 导出数据 CPK 分析',
          '功能三：立讯 MES 导出数据 CPK 分析',
          '支持批次、工站、产品型号多维度分析'],
         Inches(9.0), Inches(1.45), Inches(3.9), Inches(2.5),
         title_bg=C_ORANGE, body_bg=RGBColor(0xFF,0xF8,0xE1))

    card(slide, '技术栈',
         ['GUI：Python tkinter（5 Section 自适应布局，深色日志窗口）',
          '数据处理：pandas · openpyxl · shutil',
          '故障知识库：SQLite（fault_rules / fault_records / fault_stats）',
          '报告生成：自包含 HTML + Chart.js（CDN）',
          'LLM 增强：Ollama（localhost:11434，可选，推荐 qwen2.5:7b）'],
         Inches(0.3), Inches(4.1), Inches(12.7), Inches(2.8))


def slide_architecture(prs):
    slide = blank_slide(prs)
    title_bar(slide, '系统架构', '文件模块职责划分')
    footer(slide)

    modules = [
        ('main.py', C_RED,
         ['主窗口 + 事件驱动 GUI（LocalAnalysisTab）',
          'StationRow / MergeRuleRow 配置行组件',
          '_run_analysis() 后台线程主流程',
          'app_config.json 持久化读写',
          '菜单：工具→加载故障规则 / 帮助']),
        ('core/data_extractor.py', C_DARK,
         ['run_extraction()：latest_pass/all/fail_only 模式',
          'run_extraction_all_pass()：全部成功记录遍历',
          'run_extraction_traverse()：所选文件夹全量遍历',
          'discover_barcodes()：自动发现条码',
          'generate_missing/duplicate/folder_direct_excel()']),
        ('core/cpk_calculator.py', C_GREEN,
         ['analyze_xlsx_folder() / analyze_json_folder()',
          '读取 Test_Result_*.xlsx 或 *_MEASUREMENT_*.json',
          '计算均值、标准差、Cp、Cpk',
          '_file_time_from_name()：文件名时间戳优先',
          'values 存 (barcode, value, is_pass) 三元组']),
        ('core/fault_analyzer.py', RGBColor(0x5E,0x35,0xB1),
         ['run_fault_analysis()：结构化日志解析+规则匹配',
          'Ollama LLM 增强（基础版/增强版可选）',
          'generate_fault_barcode_list()：故障条码Excel',
          'generate_rule_suggestions_yaml()：规则建议模板',
          '跨站分析：get_cross_station_barcodes()']),
        ('core/fault_db.py', RGBColor(0x4E,0x34,0x2E),
         ['SQLite 三表：fault_rules/fault_records/fault_stats',
          '16 类种子规则预置',
          'init_db() 含自动 schema 迁移',
          'add_rule / update_rule / get_rules API']),
        ('core/html_*.py', RGBColor(0x00,0x69,0x6F),
         ['html_comprehensive_report.py：6 Tab 综合报告',
          'html_fail_report.py：帕累托+汇总卡+失败条码表',
          'html_report.py：CPK 专项报告（按工站+大项）']),
    ]

    cols = 3
    cw = Inches(4.2)
    rh = Inches(2.45)
    ml = Inches(0.2)
    mt = Inches(1.5)
    gap = Inches(0.12)

    for i, (name, color, lines) in enumerate(modules):
        col = i % cols
        row = i // cols
        l = ml + col * (cw + gap)
        t = mt + row * (rh + gap)
        card(slide, name, lines, l, t, cw, rh,
             title_bg=color, body_sz=Pt(10))


def slide_gui(prs):
    slide = blank_slide(prs)
    title_bar(slide, 'GUI 界面布局', '5 Section 自适应布局 · v2.0')
    footer(slide)

    # Mock GUI frame
    gl = Inches(0.3)
    gt = Inches(1.5)
    gw = Inches(7.1)
    gh = Inches(5.5)
    add_rect(slide, gl, gt, gw, gh, fill=C_WHITE, line=C_RED, line_w=Pt(1.5))
    # Title bar of mock GUI
    add_rect(slide, gl, gt, gw, Inches(0.3), fill=C_RED)
    add_text(slide, '产线数据分析AI平台 — Zillnk Efficiency Improvement Group',
             gl + Inches(0.1), gt + Pt(2), gw - Inches(0.15), Inches(0.28),
             size=Pt(8), color=C_WHITE)

    # Tabs
    for i, (tab, active) in enumerate([('本地数据分析', True), ('深科技MES分析', False), ('立讯MES分析', False)]):
        tl = gl + Inches(0.12) + i * Inches(1.75)
        tc = C_CREAM if active else C_RED
        add_rect(slide, tl, gt + Inches(0.3), Inches(1.65), Inches(0.27), fill=tc, line=C_RED, line_w=Pt(0.5))
        add_text(slide, tab, tl + Inches(0.05), gt + Inches(0.3),
                 Inches(1.65), Inches(0.27),
                 size=Pt(7.5), bold=active,
                 color=C_RED if active else C_WHITE, align=PP_ALIGN.CENTER)

    # Section blocks
    sec_t = gt + Inches(0.58)
    sections = [
        ('Section 1  输入/输出配置', Inches(0.38)),
        ('Section 2/3  工站配置（左）+  分析模式（右）', Inches(1.8)),
        ('Section 4  [ 开始分析 ]  [ 查看日志 ]  ████████░░ 进度', Inches(0.52)),
    ]
    ct = sec_t
    for sec_title, sec_h in sections:
        add_rect(slide, gl + Inches(0.1), ct, gw - Inches(0.2), sec_h,
                 fill=C_PINK, line=C_RED, line_w=Pt(0.75))
        add_text(slide, sec_title, gl + Inches(0.2), ct + Pt(3),
                 gw - Inches(0.35), sec_h - Pt(4),
                 size=Pt(8.5), color=C_RED)
        ct += sec_h + Inches(0.05)

    # Right panel
    dl = Inches(7.6)
    dt = Inches(1.5)

    card(slide, 'Section 1 — 输入/输出配置',
         ['DUT条码 Excel（可选，用于过滤条码）',
          '输出目录（自动创建 产品类别_时间戳 子目录）'],
         dl, dt, Inches(5.3), Inches(1.35))

    card(slide, 'Section 2 — 测试工站配置',
         ['工站类型标签 + 数据文件夹路径',
          '可折叠"工站合并配置"（源→目标合并）',
          '↑↓ 键行间切换，✕ 按钮删除行'],
         dl, dt + Inches(1.45), Inches(5.3), Inches(1.5))

    card(slide, 'Section 3 — 分析模式',
         ['xlsx / json 文件类型单选',
          '5种模式下拉（悬停显示提示气泡）',
          '故障分析方式（含失败/仅失败时显示）'],
         dl, dt + Inches(3.05), Inches(5.3), Inches(1.35))

    card(slide, 'Section 4 — 操作控制',
         ['开始分析 ↔ 停止分析（实时切换）',
          '查看运行日志（独立深色弹窗）',
          '进度标签 + 进度条实时更新'],
         dl, dt + Inches(4.5), Inches(5.3), Inches(1.2))


def slide_modes(prs):
    slide = blank_slide(prs)
    title_bar(slide, '5 种分析模式', '覆盖发货批次分析 · 全量过程分析 · 失败根因分析')
    footer(slide)

    modes = [
        ('最后一次 pass 数据', 'latest_pass', C_RED,
         '需 DUT 条码 Excel',
         '取每个条码最新一次成功记录做 CPK。\n适用：发货批次过程能力分析。\n输出：missing_barcodes.xlsx + 报告'),
        ('所选文件夹分析', 'folder_direct', C_DARK,
         '可选 Excel 过滤',
         '遍历所有配置工站目录，提取 pass+fail 全量记录。\n有失败时额外生成：fail_analysis_report.html\n+ folder_direct_fail_analysis.xlsx（3 Sheet）'),
        ('全部成功数据', 'all_pass', C_GREEN,
         '无需 Excel，自动扫描',
         '收集全部工站成功记录（含同一模块多次通过）。\n额外输出：duplicate_barcodes.xlsx（重复条码统计）'),
        ('所有数据（含失败）', 'all_with_fail', RGBColor(0x5E,0x35,0xB1),
         '无需 Excel，触发故障分析',
         'pass+fail 全量 CPK，自动建立故障分析库。\n重点：跨站比对（同模块不同设备数据关联）。\n输出：fault_barcodes.xlsx + rule_suggestions.yaml'),
        ('仅失败数据', 'fail_only', C_ORANGE,
         '无需 Excel，触发故障分析',
         '仅失败记录 CPK，深度分析失败规律。\n高频失败项 / 设备错误 Top 排行 / 未分类样本。\n输出：fault_barcodes.xlsx + rule_suggestions.yaml'),
    ]

    for i, (name, val, color, req, desc) in enumerate(modes):
        col = i % 3
        row = i // 3
        l = Inches(0.22) + col * Inches(4.36)
        t = Inches(1.5) + row * Inches(2.85)
        w = Inches(4.2)
        h = Inches(2.6)

        add_rect(slide, l, t, w, h, fill=C_CREAM, line=color, line_w=Pt(1.5))
        add_rect(slide, l, t, w, Inches(0.08), fill=color)

        add_text(slide, name, l + Inches(0.1), t + Inches(0.1), w - Inches(0.15), Inches(0.4),
                 size=Pt(13), bold=True, color=color)
        add_rect(slide, l + Inches(0.1), t + Inches(0.5), Inches(1.5), Inches(0.22), fill=color)
        add_text(slide, val, l + Inches(0.13), t + Inches(0.5), Inches(1.5), Inches(0.22),
                 size=Pt(8.5), color=C_WHITE)
        add_text(slide, req, l + Inches(1.7), t + Inches(0.51), w - Inches(1.8), Inches(0.22),
                 size=Pt(8.5), color=color, bold=True)
        add_text(slide, desc, l + Inches(0.1), t + Inches(0.82), w - Inches(0.18), Inches(1.65),
                 size=Pt(10), color=C_GRAY)


def slide_flow(prs):
    slide = blank_slide(prs)
    title_bar(slide, '数据提取与处理流程', '以「所选文件夹分析」模式为例（全量遍历）')
    footer(slide)

    steps = [
        ('① 配置工站目录', '多个工站类型\n各自配置文件夹\n（FT1/FT2/Aging…）', C_RED),
        ('② 全量遍历', '_walk_all_records\n_in_folder()\n扫描barcode/ts结构', C_DARK),
        ('③ 条码过滤', '可选：按 DUT\nExcel 筛选条码\n（保留指定批次）', RGBColor(0x00,0x69,0x6F)),
        ('④ 文件复制', 'xlsx/json 复制\n到输出目录\n按工站类型归档', C_GREEN),
        ('⑤ CPK 计算', 'analyze_xlsx\n_folder()\n计算 Cp/Cpk', RGBColor(0x5E,0x35,0xB1)),
        ('⑥ 报告生成', 'CPK 报告\n综合报告（6Tab）\n失败分析报告', C_ORANGE),
    ]

    bw = Inches(1.85)
    bh = Inches(1.55)
    aw = Inches(0.3)
    total = len(steps) * bw + (len(steps) - 1) * aw
    sl = (W - total) / 2
    ct = Inches(3.2)

    for i, (title, body, color) in enumerate(steps):
        l = sl + i * (bw + aw)
        add_rect(slide, l, ct - bh / 2, bw, bh, fill=color)
        add_text(slide, title, l + Inches(0.06), ct - bh / 2 + Inches(0.05), bw - Inches(0.1), Inches(0.36),
                 size=Pt(10.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, body, l + Inches(0.06), ct - bh / 2 + Inches(0.42), bw - Inches(0.1), Inches(1.05),
                 size=Pt(9), color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, '▶', l + bw + Inches(0.03), ct - Inches(0.15), aw, Inches(0.3),
                     size=Pt(15), color=C_RED, align=PP_ALIGN.CENTER)

    notes = [
        '目录结构：<工站根> / TestResult / <产品> / <工站类型> / [中间层] / <条码> / <时间戳> / xlsx + json',
        '条码文件夹识别：直接子目录含时间戳格式（YYYYMMDDHHMMSS）文件夹，即视为条码文件夹',
        '跳过规则：debug / file_bk / env_comp / RU*_Log_* 等目录自动跳过，深度上限 10 层',
        '工站合并：可将多个工站类型（如 FT1 + FT2）合并为同一 CPK 分析组',
    ]
    nt = ct + bh / 2 + Inches(0.3)
    for j, note in enumerate(notes):
        add_text(slide, f'  ◆  {note}', Inches(0.35), nt + j * Inches(0.4),
                 Inches(12.6), Inches(0.38), size=Pt(9.5), color=C_GRAY)


def slide_fault(prs):
    slide = blank_slide(prs)
    title_bar(slide, '故障分析模块', '规则库匹配 + Ollama LLM 增强 · 持续积累知识')
    footer(slide)

    card(slide, '基础版（规则库）',
         ['内置 16 类种子规则，覆盖主要故障类型',
          '优先级：设备/通信错误 > 失败测试项名 > 日志关键词',
          'CRITICAL 日志行结构化解析：测试项+测量值+限值+偏差',
          'env_config.yml → VISA 地址映射 → EQP_ID 设备标识',
          '设备错误 6 种正则：COM断连/VISA/串口超时/SSH/网络/射频开关'],
         Inches(0.3), Inches(1.5), Inches(6.15), Inches(2.7))

    card(slide, '增强版（规则库 + Ollama）',
         ['在规则库基础上，调用本地 Ollama LLM 辅助分析',
          '端点：localhost:11434（推荐模型 qwen2.5:7b）',
          '对未分类故障给出根因推断和处置建议',
          '分析结果写入 fault_database.db（可持续积累）',
          '仅对"所有数据（含失败）"和"仅失败数据"模式触发'],
         Inches(6.6), Inches(1.5), Inches(6.4), Inches(2.7),
         title_bg=C_DARK)

    card(slide, '知识库持续积累流程',
         ['1. 分析完成 → 生成 rule_suggestions_<时间戳>.yaml（含未分类样本 + Top10 失败项）',
          '2. 工程师用文本编辑器填写 fault_type / suggestion（空白条目）',
          '3. 菜单 → 工具 → 加载故障关系描述文件… → 选择 YAML',
          '4. 系统自动合并入 fault_database.db（相同关键词更新，不重复导入）',
          '5. 下次分析时新规则立即生效，未分类故障比例持续下降'],
         Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.65),
         title_bg=C_GREEN)


def slide_reports(prs):
    slide = blank_slide(prs)
    title_bar(slide, '输出文件 & HTML 报告', '每次运行隔离存放在 产品类别_时间戳/ 子目录')
    footer(slide)

    card(slide, '输出文件一览（按模式）',
         ['missing_barcodes.xlsx       — 缺失/异常条码明细（latest_pass / all_with_fail / fail_only）',
          'duplicate_barcodes.xlsx     — 重复测试条码统计（all_pass）',
          'fault_barcodes.xlsx         — 故障条码分类列表（all_with_fail / fail_only）',
          'folder_direct_fail_analysis.xlsx  — 失败明细 3Sheet（所选文件夹分析）',
          'rule_suggestions_<ts>.yaml  — 规则建议模板，供工程师填写',
          'cpk_report.html             — CPK 专项报告（按工站+测试大项分组，所有模式）',
          'comprehensive_report.html   — 综合分析报告 6 Tab（所有模式，自动浏览器打开）',
          'fail_analysis_report.html   — 帕累托+汇总卡+失败条码表（所选文件夹分析）',
          'fault_database.db           — 故障关系知识库（SQLite）',
          'analysis_log_<ts>.txt       — 本次运行完整日志'],
         Inches(0.3), Inches(1.5), Inches(7.9), Inches(5.65),
         body_sz=Pt(10.5))

    card(slide, '综合报告（comprehensive_report.html）6 个 Tab',
         ['总览        — KPI 卡片 + 良率趋势 + 失败类型分布 + 测试大类汇总',
          '失败分析    — Top25 高频失败项柱图（点击查看明细）+ 失败记录表',
          'CPK 分析    — Cpk 横向柱图（色标分级：<1红/1~1.33橙/≥1.33绿）+ 完整 CPK 表',
          '数据分布    — 按测试项切换堆叠直方图（pass蓝/fail红）+ 统计面板',
          '失败模式    — 失败类型统计 + 时序热图 + 多失败项 SN 关联分析',
          '故障回放    — 左侧 SN 列表（可搜索/过滤）+ 右侧逐项测试结果展开'],
         Inches(8.4), Inches(1.5), Inches(4.6), Inches(5.65),
         title_bg=RGBColor(0x00,0x69,0x6F), body_sz=Pt(10.5))


def slide_roadmap(prs):
    slide = blank_slide(prs)
    title_bar(slide, '开发路线图', 'P1 优先级 → P2 → 后期功能')
    footer(slide)

    card(slide, '✅ 已完成（v2.0）',
         ['5种分析模式全量提取 + CPK 计算',
          '综合 HTML 报告（6 Tab，Chart.js）',
          '失败分析报告（帕累托 + 汇总卡）',
          '故障关系知识库（SQLite，16 类种子规则）',
          'Ollama LLM 增强分析（可选）',
          'YAML 规则导入工具（菜单入口）',
          '工站合并配置（持久化）',
          'GUI 自适应高度 + 深色日志弹窗'],
         Inches(0.3), Inches(1.5), Inches(3.85), Inches(5.65),
         title_bg=C_GREEN, body_sz=Pt(10.5))

    card(slide, '🔧 P1 — 近期优先',
         ['HTML 报告重构：两级 Tab / 隐藏 CPK 列 / 全宽分布图 / 产品数据检索 Tab',
          'env_comp/*.csv 仪器校准漂移检测',
          'test_limits.yml 跨批次测试限值哈希比对（限值变更告警）'],
         Inches(4.35), Inches(1.5), Inches(4.2), Inches(2.6),
         title_bg=C_ORANGE, body_sz=Pt(11))

    card(slide, '📋 P2 — 中期规划',
         ['Pass/Fail 差异比对引擎（barcode_comparisons 表）',
          'load_knowledge_yaml() 后端独立封装',
          '时间窗口聚类仪器告警',
          '共失败关联性矩阵',
          'all_with_fail 跨站推断自动化'],
         Inches(4.35), Inches(4.3), Inches(4.2), Inches(2.85),
         title_bg=C_DARK, body_sz=Pt(11))

    card(slide, '🚀 后期功能',
         ['功能二：深科技 MES 导出数据 CPK 分析',
          '（支持批次 / 工站 / 产品型号多维度）',
          '功能三：立讯 MES 导出数据 CPK 分析',
          '（支持批次 / 工站 / 产品型号多维度）'],
         Inches(8.75), Inches(1.5), Inches(4.25), Inches(5.65),
         title_bg=C_RED, body_sz=Pt(11))


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = init_prs()

    slide_title(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_gui(prs)
    slide_modes(prs)
    slide_flow(prs)
    slide_fault(prs)
    slide_reports(prs)
    slide_roadmap(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       '产线数据分析AI平台_架构说明.pptx')
    prs.save(out)
    print(f'已生成: {out}')
    print(f'共 {len(prs.slides)} 张幻灯片（公司模板 · 红白配色 · 含Logo）')


if __name__ == '__main__':
    main()
